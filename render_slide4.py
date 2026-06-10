"""Faithful-ish PIL preview of slide 4 (table + text) read from the saved pptx."""
import os
import matplotlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

DPI = 150.0
EMU = 914400.0
_FT = os.path.join(os.path.dirname(matplotlib.__file__), "mpl-data", "fonts", "ttf")


def font(sz, bold=True):
    p = os.path.join(_FT, "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf")
    try:
        return ImageFont.truetype(p, max(8, int(sz)))
    except Exception:
        return ImageFont.load_default()


def px(v):
    return int(v / EMU * DPI)


prs = Presentation("10_Autopass_Gen.pptx")
W = px(prs.slide_width); H = px(prs.slide_height)
img = Image.new("RGB", (W, H), "white")
d = ImageDraw.Draw(img)
s = prs.slides[3]


def rgb(c, default=(0, 0, 0)):
    try:
        return (c[0], c[1], c[2])
    except Exception:
        return default


for sh in s.shapes:
    L, T, Wd, Hd = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
    if sh.has_table:
        t = sh.table
        nC = len(t.columns); nR = len(t.rows)
        colw = [px(t.columns[c].width) for c in range(nC)]
        rowh = [px(t.rows[r].height) for r in range(nR)]
        # scale row heights to fit Hd
        sh_total = sum(rowh)
        rowh = [int(h * Hd / sh_total) for h in rowh]
        y = T
        for r in range(nR):
            x = L
            for c in range(nC):
                cell = t.cell(r, c)
                fill = (255, 255, 255)
                try:
                    fill = rgb(cell.fill.fore_color.rgb, (255, 255, 255))
                except Exception:
                    pass
                d.rectangle([x, y, x + colw[c], y + rowh[r]], fill=fill, outline=(200, 206, 212), width=1)
                # text
                paras = [p for p in cell.text_frame.paragraphs if p.text.strip()]
                line_h = rowh[r] / max(1, len(paras) + 0.4)
                ty = y + (rowh[r] - line_h * len(paras)) / 2
                for p in paras:
                    run = p.runs[0] if p.runs else None
                    sz = run.font.size.pt if (run and run.font.size) else 14
                    col = rgb(run.font.color.rgb, (0, 0, 0)) if run and run.font.color and run.font.color.type is not None else (0, 0, 0)
                    b = bool(run.font.bold) if run else False
                    f = font(sz / 72.0 * DPI, b)
                    txt = p.text
                    tw = d.textlength(txt, font=f)
                    align = str(p.alignment)
                    tx = x + 8 if "LEFT" in align else x + (colw[c] - tw) / 2
                    d.text((tx, ty), txt, font=f, fill=col)
                    ty += line_h
                x += colw[c]
            y += rowh[r]
    elif sh.has_text_frame and sh.text_frame.text.strip():
        y = T
        for p in sh.text_frame.paragraphs:
            if not p.text.strip():
                continue
            run = p.runs[0] if p.runs else None
            sz = run.font.size.pt if (run and run.font.size) else 16
            col = rgb(run.font.color.rgb, (0, 0, 0)) if run and run.font.color and run.font.color.type is not None else (0, 0, 0)
            b = bool(run.font.bold) if run else False
            f = font(sz / 72.0 * DPI, b)
            tw = d.textlength(p.text, font=f)
            align = str(p.alignment)
            tx = L + 6 if "LEFT" in align or align == "None" else L + (Wd - tw) / 2
            d.text((tx, y), p.text, font=f, fill=col)
            y += int(f.size * 1.3)

img.save("presentation_assets/preview_slide4.png")
print("wrote presentation_assets/preview_slide4.png", img.size)
