from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
try:
    from pptx.enum.text import MSO_ANCHOR
except Exception:
    MSO_ANCHOR = None

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
INK = RGBColor(0x16, 0x22, 0x2E)
GREEN = RGBColor(0x1F, 0x82, 0x3C)
GRAY = RGBColor(0x5B, 0x6B, 0x7B)
RED = RGBColor(0xC0, 0x3A, 0x2E)
AMBER = RGBColor(0xB5, 0x6A, 0x16)
WHITEc = RGBColor(0xFF, 0xFF, 0xFF)
OURS = RGBColor(0xE6, 0xF2, 0xEA)   # highlight row for AutoPass-Gen

prs = Presentation("00_ProjectTitle.pptx")
slides = list(prs.slides)


def by_id(slide, sid):
    for s in slide.shapes:
        if s.shape_id == sid:
            return s
    raise KeyError(sid)


def set_run_text(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = text
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)


def set_paras(shape, items, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    tf.paragraphs[0].clear()
    first = True
    for txt, bold, size, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"


def replace_with_image(slide, shape, img):
    iw, ih = Image.open(img).size
    L, T, W, H = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    sc = min(W / iw, H / ih)
    w, h = int(iw * sc), int(ih * sc)
    slide.shapes.add_picture(img, L + (W - w) // 2, T + (H - h) // 2, width=w, height=h)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _cell(cell, text, size, color, fill, bold=True, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    if MSO_ANCHOR is not None:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    cell.margin_left = Pt(8); cell.margin_right = Pt(8)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, part in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = part
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = "Calibri"; r.font.color.rgb = color


def add_comparison_table(slide, left, top, width):
    """Native, editable, colour-coded 3-policy comparison table from the CARLA benchmark."""
    rows, cols = 4, 6
    gfx = slide.shapes.add_table(rows, cols, left, top, width, Inches(2.9))
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, w in enumerate([3.1, 1.3, 1.15, 1.45, 1.25, 0.95]):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(0.66)
    for i in (1, 2, 3):
        tbl.rows[i].height = Inches(0.7)
    headers = ["Policy", "Overtakes\ncompleted", "Collisions",
               "Unsafe\npass attempts", "Mean speed\n(m/s)", "Lane dev\n(m)"]
    for c, h in enumerate(headers):
        _cell(tbl.cell(0, c), h, 13, WHITEc, NAVY, bold=True)

    def row(i, name, data, fill, name_color=INK):
        _cell(tbl.cell(i, 0), name, 14, name_color, fill, bold=True, align=PP_ALIGN.LEFT)
        for c, (txt, col) in enumerate(data, start=1):
            _cell(tbl.cell(i, c), txt, 16, col, fill, bold=True)

    row(1, "No-pass  (never overtake)",
        [("0 / 18", AMBER), ("0 / 23", GREEN), ("0", GREEN), ("3.9", AMBER), ("0.24", INK)], WHITEc)
    row(2, "AutoPass-Gen  (ours)",
        [("18 / 18", GREEN), ("0 / 23", GREEN), ("0", GREEN), ("9.7", GREEN), ("0.60", INK)], OURS, name_color=NAVY)
    row(3, "Aggressive  (always pass)",
        [("18 / 18", GREEN), ("2 / 23", RED), ("2", RED), ("9.7", GREEN), ("0.60", INK)], WHITEc)
    return gfx


# ===================== SLIDE 1 - Title =====================
s = slides[0]
bar = by_id(s, 8)  # remove the empty grey accent bar (looks unfinished)
bar._element.getparent().remove(bar._element)
set_run_text(by_id(s, 6), "AutoPass-Gen")
set_paras(by_id(s, 2052),
          [("Vision-grounded agentic overtaking under trip-deadline pressure", False, 18, NAVY)],
          align=PP_ALIGN.CENTER)
set_paras(by_id(s, 2050),
          [("Team 10      Rohan Sachdeva  ·  Xinwei Mai  ·  Pranav Prabu  ·  Rathang Pandit      UC San Diego, CSE 252D",
            False, 14, INK)],
          align=PP_ALIGN.CENTER)
set_paras(by_id(s, 13),
          [("An LLM agent chooses which perception tools to run, a deterministic critic gates every overtake, "
            "and a waypoint controller keeps the car in its lane, so it passes when it is safe and waits when it is not.",
            False, 16, INK)],
          align=PP_ALIGN.CENTER)
replace_with_image(s, by_id(s, 3), "presentation_assets/fig_title_hero.png")
set_notes(s,
    "[ TOTAL TARGET 4:00 - suggested split: Rohan = slides 1,2,4,5 ; Xinwei = slide 3 (the agent). Practice once out loud. ]\n\n"
    "SLIDE 1 - TITLE  (~18s, Rohan)\n"
    "\"We're Team 10. Our project, AutoPass-Gen, is an autonomous car that decides when to overtake a slow "
    "vehicle under a trip deadline - using only what it can see - and it makes that decision as an agent that "
    "picks its own tools and verifies its own safety. Let me set up why that's hard.\"")

# ===================== SLIDE 2 - Problem =====================
s = slides[1]
replace_with_image(s, by_id(s, 3), "presentation_assets/fig_problem.png")
set_paras(by_id(s, 7),
          [("Overtake a slow lead, and still make the deadline.", True, 20, NAVY),
           ("A safe pass depends on the rear and oncoming lanes and on timing, not on a pre-planned trajectory. "
            "A wrong call means a missed deadline, or a crash.", False, 16, INK)])
set_notes(s,
    "SLIDE 2 - PROBLEM  (~45s, Rohan)\n"
    "\"The maneuver is overtaking. Staying behind a slow car is always safe - but you can miss your deadline. "
    "Passing makes progress - but it can be deadly. The key point is this: whether a pass is safe is NOT a "
    "pre-planned path. It depends on what's happening around the car right now - is a fast car closing in the "
    "passing lane? is there oncoming traffic? is the time-gap big enough? A trajectory predictor that just "
    "extrapolates motion can't answer 'should I pull out right now' - it has to actively check the other lanes. "
    "So the decision we target is simply: pass, or wait - and only pass when it's truly safe.\"")

# ===================== SLIDE 3 - Unique Perspective =====================
s = slides[2]
replace_with_image(s, by_id(s, 3), "presentation_assets/fig_architecture.png")
set_paras(by_id(s, 7),
          [("The agent controls the process, not just the wording.", True, 20, NAVY),
           ("LLM planner chooses tools each cycle  →  deterministic critic verifies safety  →  re-plan on reject. "
            "Greedy under urgency, safe under code.", False, 16, INK)])
set_notes(s,
    "SLIDE 3 - UNIQUE PERSPECTIVE / THE AGENT  (~75s, Xinwei)\n"
    "\"Here's why this is agentic, not a single model. Every half-second the car runs a loop. A planner - an "
    "LLM - chooses which tool to call next: look at the front gap, the rear lane, the oncoming lane, check the "
    "road geometry, or propose a pass. It is NOT a fixed script - the order and the number of tool calls change "
    "with the scene, and everything it learns is written into a shared, mutable state that also remembers past "
    "rejections. When the planner proposes a pass, it does not get to approve itself: a separate, deterministic "
    "critic checks the hard safety gates - front gap, rear time-to-collision, oncoming clearance. If the critic "
    "says unsafe, the agent goes back, senses again, and re-plans. If it says safe, only then does a waypoint "
    "controller run the lane change - and because it steers to real lane-center waypoints, the car stays in its "
    "lane by construction. The one insight that ties it together: the LLM owns the PROCESS, the critic owns "
    "SAFETY. Greedy under urgency, but safe under code.\"")

# ===================== SLIDE 4 - Results =====================
s = slides[3]
# The quantitative comparison table is the centerpiece (native + editable).
_fig = by_id(s, 3)
_fig._element.getparent().remove(_fig._element)
_cap = s.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.4))
set_paras(_cap, [("Same 23 CARLA scenarios across 4 towns, three decision policies (live agentic autopass)",
                  False, 13, GRAY)], align=PP_ALIGN.CENTER)
add_comparison_table(s, Inches(0.4), Inches(1.45), Inches(9.2))
set_paras(by_id(s, 7),
          [("Full campaign: 23 / 23 correct, 0 collisions, 4 towns.", True, 20, NAVY),
           ("What worked: the critic + re-plan loop removes the aggressive baseline's crashes, "
            "while urgency removes the no-pass baseline's stalling.", False, 16, INK)])
set_notes(s,
    "SLIDE 4 - AMAZING RESULTS  (~50s, Rohan)\n"
    "\"Now the numbers. We compared three policies on the SAME CARLA scenarios. A never-pass car is safe but "
    "never overtakes and crawls at about a third of the speed. An always-pass car is fast, but it crashes - "
    "two collisions and three unsafe attempts in six scenarios - because it ignores the rear and oncoming "
    "lanes. AutoPass-Gen gets every overtake at full speed with zero collisions and zero unsafe passes. And "
    "across the full campaign - 23 of 23 scenarios in four towns - it was correct every time. The reason is "
    "the agentic part: the deterministic critic vetoes unsafe passes and the agent re-plans, while urgency "
    "keeps it from stalling like the never-pass baseline.\"")

# ===================== SLIDE 5 - Demo =====================
s = slides[4]
vshape = by_id(s, 3)
L, T, W, H = vshape.left, vshape.top, vshape.width, vshape.height
vshape._element.getparent().remove(vshape._element)
iw, ih = Image.open("presentation_assets/demo_poster.png").size
vw = int(W)
vh = int(W * ih / iw)
vleft = L
vtop = T + Emu(int(0.25 * 914400))
s.shapes.add_movie("presentation_assets/demo_hero.mp4", vleft, vtop, vw, vh,
                   poster_frame_image="presentation_assets/demo_poster.png", mime_type="video/mp4")
cap_top = vtop + vh + Emu(int(0.28 * 914400))
tb = s.shapes.add_textbox(L, cap_top, W, Emu(int(1.5 * 914400)))
set_paras(tb,
          [("Two cases a trajectory predictor cannot do: read the next lane and decide to WAIT.", True, 20, NAVY),
           ("1)  Highway: yield to a fast car closing in the passing lane, then overtake.", False, 16, INK),
           ("2)  Rural two-lane: wait for oncoming traffic to clear, then overtake.", False, 16, INK)],
          align=PP_ALIGN.CENTER)
set_notes(s,
    "SLIDE 5 - AMAZING DEMO  (~50s incl. video, Rohan)\n"
    "\"Let me show you the two cases that make the point - watch the lane NEXT to the car.\"  [ CLICK THE VIDEO TO PLAY ]\n"
    "(while it plays) \"First, a highway: the car wants to pass, but a fast vehicle is closing in the passing "
    "lane - watch it hold, let that car go by, and only then overtake. ... Now a rural two-lane: the only "
    "passing lane is the oncoming lane - it waits for the oncoming car to clear, then completes the pass.\"\n"
    "(after) \"Neither of these is replaying a trajectory - the car is reading the lane next to it and deciding "
    "to wait. That's our whole thesis, working end to end. Thank you.\"\n\n"
    "[ BACKUP: if the embedded video fails, play 10_Autopass_Gen.mp4 from the same folder. ]")

import time as _time
for _cand in ["10_Autopass_Gen_rev.pptx", f"10_Autopass_Gen_rev_{int(_time.time())}.pptx"]:
    try:
        prs.save(_cand)
        print(f"saved {_cand} ; slides:", len(prs.slides._sldIdLst))
        break
    except PermissionError:
        print(f"LOCKED (open in PowerPoint?): {_cand}")
